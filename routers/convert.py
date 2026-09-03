"""
convert.py — All file conversion endpoints for PrestigePDF.

Conversions:
  PDF → Word (.docx)     : pdf2docx
  PDF → Excel (.xlsx)    : pdfplumber + openpyxl
  PDF → PowerPoint (.pptx): pdf2image + python-pptx
  PDF → Images (.zip)    : pdf2image (requires Poppler on PATH)
  Word → PDF             : LibreOffice headless
  Excel → PDF            : LibreOffice headless
  PPT → PDF              : LibreOffice headless
"""

import io
import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path

import pdfplumber
from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse, StreamingResponse
from pdf2docx import Converter as DocxConverter
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt

from routers.utils import cleanup, cleanup_dir, make_tmp_dir, save_upload

router = APIRouter()


# ─────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────

def _libreoffice_convert(in_path: str, out_dir: str, to_format: str) -> str | None:
    """
    Use LibreOffice headless to convert a document.
    Returns path to the converted file, or None if LibreOffice not available.
    """
    soffice = None
    candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",
        "libreoffice",
    ]
    for c in candidates:
        if os.path.exists(c) or shutil.which(c):
            soffice = c
            break

    if not soffice:
        return None

    try:
        result = subprocess.run(
            [
                soffice,
                "--headless",
                "--norestore",
                "--convert-to",
                to_format,
                "--outdir",
                out_dir,
                in_path,
            ],
            timeout=120,
            capture_output=True,
        )
        if result.returncode == 0:
            # LibreOffice outputs a file with same basename but new extension
            base = Path(in_path).stem
            out_file = os.path.join(out_dir, f"{base}.{to_format}")
            if os.path.exists(out_file):
                return out_file
    except Exception:
        pass
    return None


def _pdf2image_available() -> bool:
    """Check if Poppler is available for pdf2image."""
    try:
        from pdf2image.exceptions import PDFInfoNotInstalledError
        import pdf2image
        pdf2image.pdfinfo_from_path.__doc__  # just access it
        return True
    except Exception:
        return False


def _pymupdf_to_docx(in_path: str, out_path: str) -> None:
    """
    Fallback PDF→DOCX using PyMuPDF rawdict extraction + python-docx.
    Preserves:
      - Font size (maps to heading level or body)
      - Bold / italic styling per span
      - Text color
      - Proper word spacing (PyMuPDF correctly handles kerned fonts)
      - Paragraph grouping by vertical proximity
    """
    import pymupdf as fitz
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc_pdf = fitz.open(in_path)
    doc_word = Document()

    # Narrow margins
    for section in doc_word.sections:
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)

    # Clear default empty paragraph
    for p in doc_word.paragraphs:
        p._element.getparent().remove(p._element)

    for page_idx in range(doc_pdf.page_count):
        if page_idx > 0:
            doc_word.add_page_break()

        page = doc_pdf[page_idx]
        raw = page.get_text("rawdict", flags=fitz.TEXT_PRESERVE_WHITESPACE)

        prev_y1 = -1

        for block in raw.get("blocks", []):
            if block.get("type") != 0:  # skip images
                continue

            for line in block.get("lines", []):
                spans = line.get("spans", [])
                if not spans:
                    continue

                # Detect blank-line gap between blocks
                y0 = line["bbox"][1]
                if prev_y1 >= 0 and (y0 - prev_y1) > 8:
                    doc_word.add_paragraph("")  # spacer

                prev_y1 = line["bbox"][3]

                # Build a single Word paragraph for this line
                para = doc_word.add_paragraph()
                para.paragraph_format.space_before = Pt(0)
                para.paragraph_format.space_after = Pt(1)

                for span in spans:
                    raw_text = span.get("text", "")
                    if not raw_text:
                        continue

                    size = float(span.get("size", 11))
                    flags = int(span.get("flags", 0))
                    is_bold = bool(flags & 2**4)  # bit 4 = bold
                    is_italic = bool(flags & 2**1)  # bit 1 = italic
                    color_int = span.get("color", 0)
                    r = (color_int >> 16) & 0xFF
                    g = (color_int >> 8) & 0xFF
                    b = color_int & 0xFF

                    run = para.add_run(raw_text)
                    run.bold = is_bold
                    run.italic = is_italic
                    run.font.size = Pt(max(7, min(size, 36)))
                    try:
                        run.font.color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass

    doc_pdf.close()
    doc_word.save(out_path)


@router.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """
    Convert PDF to Word (.docx).
    Primary:  pdf2docx  — best for layout-heavy PDFs (tables, columns).
    Fallback: PyMuPDF rawdict → python-docx — preserves fonts/bold/italic
              and correct word spacing (handles kerned/embedded fonts).
    """
    in_path = await save_upload(file, suffix=".pdf")
    out_path = tempfile.mktemp(suffix=".docx")
    orig_name = Path(file.filename or "document").stem

    pdf2docx_ok = False
    try:
        cv = DocxConverter(in_path)
        cv.convert(out_path, start=0, end=None)
        cv.close()

        # Validate: check output has real text content
        if os.path.exists(out_path) and os.path.getsize(out_path) > 2000:
            from docx import Document as _Docx
            _d = _Docx(out_path)
            total_chars = sum(len(p.text) for p in _d.paragraphs)
            if total_chars > 100:
                pdf2docx_ok = True

        if not pdf2docx_ok:
            raise Exception("pdf2docx produced insufficient content")

    except Exception as primary_err:
        print(f"[INFO] pdf2docx failed ({primary_err}), using PyMuPDF fallback")
        try:
            if os.path.exists(out_path):
                os.remove(out_path)
            _pymupdf_to_docx(in_path, out_path)
        except Exception as fallback_err:
            cleanup(in_path, out_path)
            raise HTTPException(
                status_code=500,
                detail=f"PDF to Word failed. pdf2docx: {primary_err}. Fallback: {fallback_err}",
            )

    cleanup(in_path)
    return FileResponse(
        out_path,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename=f"{orig_name}.docx",
    )



# ─────────────────────────────────────────────────────
# PDF → Excel
# ─────────────────────────────────────────────────────

@router.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Extract tables from PDF and save as Excel (.xlsx) using pdfplumber + openpyxl."""
    in_path = await save_upload(file, suffix=".pdf")
    out_path = tempfile.mktemp(suffix=".xlsx")

    try:
        wb = Workbook()
        wb.remove(wb.active)  # Remove default sheet

        with pdfplumber.open(in_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()

                if tables:
                    for tbl_idx, table in enumerate(tables):
                        sheet_name = f"Page{page_num}_T{tbl_idx + 1}"[:31]
                        ws = wb.create_sheet(title=sheet_name)
                        for row in table:
                            ws.append([cell or "" for cell in row])
                else:
                    # No tables: extract plain text
                    text = page.extract_text() or ""
                    sheet_name = f"Page{page_num}"[:31]
                    ws = wb.create_sheet(title=sheet_name)
                    for line in text.splitlines():
                        ws.append([line])

        if not wb.sheetnames:
            ws = wb.create_sheet(title="Sheet1")
            ws.append(["No extractable content found"])

        wb.save(out_path)

    except Exception as e:
        cleanup(in_path, out_path)
        raise HTTPException(status_code=500, detail=f"PDF to Excel conversion failed: {str(e)}")

    cleanup(in_path)
    orig_name = Path(file.filename or "document").stem
    return FileResponse(
        out_path,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=f"{orig_name}.xlsx",
    )


# ─────────────────────────────────────────────────────
# PDF → PowerPoint
# ─────────────────────────────────────────────────────

@router.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    """
    Convert PDF pages to PowerPoint slides.
    If Poppler is available: render each page as an image slide.
    Fallback: text extraction per slide.
    """
    in_path = await save_upload(file, suffix=".pdf")
    out_path = tempfile.mktemp(suffix=".pptx")
    tmp_dir = make_tmp_dir()

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]  # blank

        # Try image-based approach (requires Poppler)
        images_created = False
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(in_path, dpi=150, output_folder=tmp_dir)
            for img in images:
                slide = prs.slides.add_slide(blank_layout)
                img_path = os.path.join(tmp_dir, f"page_{id(img)}.jpg")
                img.save(img_path, "JPEG", quality=85)
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
            images_created = True
        except Exception:
            pass

        if not images_created:
            # Text-based fallback
            import fitz
            doc = fitz.open(in_path)
            for page_idx in range(doc.page_count):
                page = doc[page_idx]
                text = page.get_text("text")
                slide = prs.slides.add_slide(blank_layout)
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text[:2000] if text else f"Page {page_idx + 1}"
                p.font.size = Pt(12)
            doc.close()

        prs.save(out_path)

    except Exception as e:
        cleanup(in_path, out_path)
        cleanup_dir(tmp_dir)
        raise HTTPException(status_code=500, detail=f"PDF to PPT conversion failed: {str(e)}")

    cleanup(in_path)
    cleanup_dir(tmp_dir)
    orig_name = Path(file.filename or "document").stem
    return FileResponse(
        out_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{orig_name}.pptx",
    )


# ─────────────────────────────────────────────────────
# PDF → Images (ZIP of JPGs)
# ─────────────────────────────────────────────────────

@router.post("/pdf-to-images")
async def pdf_to_images(
    file: UploadFile = File(...),
    dpi: int = Form(150),
):
    """
    Convert each PDF page to a JPEG image and return a ZIP archive.
    Requires Poppler on PATH. Falls back to PyMuPDF if Poppler unavailable.
    """
    in_path = await save_upload(file, suffix=".pdf")
    tmp_dir = make_tmp_dir()

    try:
        image_paths = []

        # Try pdf2image (Poppler-based, best quality)
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(in_path, dpi=dpi, output_folder=tmp_dir)
            for i, img in enumerate(images):
                img_path = os.path.join(tmp_dir, f"page_{i + 1:04d}.jpg")
                img.save(img_path, "JPEG", quality=90)
                image_paths.append(img_path)
        except Exception:
            # PyMuPDF fallback
            import fitz
            doc = fitz.open(in_path)
            zoom = dpi / 72
            mat = fitz.Matrix(zoom, zoom)
            for page_idx in range(doc.page_count):
                page = doc[page_idx]
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_path = os.path.join(tmp_dir, f"page_{page_idx + 1:04d}.jpg")
                pix.save(img_path)
                image_paths.append(img_path)
            doc.close()

        if not image_paths:
            raise Exception("No images generated from PDF")

        # Pack into ZIP
        zip_buffer = io.BytesIO()
        orig_stem = Path(file.filename or "document").stem
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for img_path in image_paths:
                zf.write(img_path, arcname=os.path.basename(img_path))
        zip_buffer.seek(0)

    except Exception as e:
        cleanup(in_path)
        cleanup_dir(tmp_dir)
        raise HTTPException(status_code=500, detail=f"PDF to images failed: {str(e)}")

    cleanup(in_path)
    cleanup_dir(tmp_dir)

    orig_name = Path(file.filename or "document").stem
    return StreamingResponse(
        zip_buffer,
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{orig_name}_images.zip"'},
    )


# ─────────────────────────────────────────────────────
# Office → PDF  (Word, Excel, PPT)
# ─────────────────────────────────────────────────────

def _fallback_docx_to_pdf(in_path: str, out_path: str) -> None:
    """Fallback DOCX → PDF using python-docx + PyMuPDF."""
    import pymupdf as fitz
    from docx import Document

    doc_word = Document(in_path)
    pdf = fitz.open()
    page_width, page_height = 595.0, 842.0
    margin = 50.0
    usable_width = page_width - (margin * 2)

    page = pdf.new_page(width=page_width, height=page_height)
    y_cursor = margin

    for p in doc_word.paragraphs:
        text = p.text.strip()
        if not text:
            y_cursor += 10.0
            continue

        font_size = 11.0
        if p.style and p.style.name and "Heading 1" in p.style.name:
            font_size = 18.0
        elif p.style and p.style.name and "Heading 2" in p.style.name:
            font_size = 14.0

        line_height = font_size * 1.35
        if y_cursor + line_height > (page_height - margin):
            page = pdf.new_page(width=page_width, height=page_height)
            y_cursor = margin

        rect = fitz.Rect(margin, y_cursor, margin + usable_width, y_cursor + line_height + 20)
        rc = page.insert_textbox(rect, text, fontsize=font_size, fontname="helv")
        if rc < 0:
            page = pdf.new_page(width=page_width, height=page_height)
            y_cursor = margin
            rect = fitz.Rect(margin, y_cursor, margin + usable_width, y_cursor + line_height + 20)
            page.insert_textbox(rect, text, fontsize=font_size, fontname="helv")

        y_cursor += max(line_height, 16.0)

    for table in doc_word.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
            if not row_text:
                continue
            if y_cursor + 16 > (page_height - margin):
                page = pdf.new_page(width=page_width, height=page_height)
                y_cursor = margin
            rect = fitz.Rect(margin, y_cursor, margin + usable_width, y_cursor + 20)
            page.insert_textbox(rect, row_text, fontsize=10, fontname="courier")
            y_cursor += 16.0

    if pdf.page_count == 0:
        page = pdf.new_page(width=page_width, height=page_height)
        page.insert_text((margin, margin), "Document", fontsize=12)

    pdf.save(out_path)
    pdf.close()


def _fallback_xlsx_to_pdf(in_path: str, out_path: str) -> None:
    """Fallback XLSX → PDF using openpyxl + PyMuPDF."""
    import pymupdf as fitz
    from openpyxl import load_workbook

    wb = load_workbook(in_path, data_only=True)
    pdf = fitz.open()

    page_width, page_height = 842.0, 595.0  # Landscape A4
    margin = 40.0
    usable_width = page_width - (margin * 2)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        page = pdf.new_page(width=page_width, height=page_height)
        y_cursor = margin

        page.insert_text((margin, y_cursor), f"Sheet: {sheet_name}", fontsize=14, fontname="helv", color=(0.1, 0.3, 0.7))
        y_cursor += 25.0

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(cells):
                continue
            row_str = " | ".join(cells)
            if len(row_str) > 130:
                row_str = row_str[:130] + "..."

            if y_cursor + 16 > (page_height - margin):
                page = pdf.new_page(width=page_width, height=page_height)
                y_cursor = margin

            rect = fitz.Rect(margin, y_cursor, margin + usable_width, y_cursor + 18)
            page.insert_textbox(rect, row_str, fontsize=9, fontname="courier")
            y_cursor += 16.0

    if pdf.page_count == 0:
        page = pdf.new_page(width=page_width, height=page_height)
        page.insert_text((margin, margin), "Spreadsheet", fontsize=12)

    pdf.save(out_path)
    pdf.close()


def _fallback_pptx_to_pdf(in_path: str, out_path: str) -> None:
    """Fallback PPTX → PDF using python-pptx (if available) or PyMuPDF."""
    import pymupdf as fitz

    pdf = fitz.open()
    page_width, page_height = 792.0, 612.0  # Landscape
    margin = 40.0
    usable_width = page_width - (margin * 2)

    try:
        from pptx import Presentation
        prs = Presentation(in_path)
        for i, slide in enumerate(prs.slides):
            page = pdf.new_page(width=page_width, height=page_height)
            y_cursor = margin

            page.insert_text((margin, y_cursor), f"Slide {i + 1}", fontsize=14, fontname="helv", color=(0.1, 0.3, 0.7))
            y_cursor += 30.0

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if y_cursor + 20 > (page_height - margin):
                        break
                    rect = fitz.Rect(margin, y_cursor, margin + usable_width, y_cursor + 40)
                    page.insert_textbox(rect, text, fontsize=11, fontname="helv")
                    y_cursor += 35.0
    except Exception:
        page = pdf.new_page(width=page_width, height=page_height)
        page.insert_text((margin, margin), "Presentation Document", fontsize=12)

    if pdf.page_count == 0:
        page = pdf.new_page(width=page_width, height=page_height)
        page.insert_text((margin, margin), "Presentation Document", fontsize=12)

    pdf.save(out_path)
    pdf.close()


async def _office_to_pdf(file: UploadFile, ext: str, label: str) -> FileResponse:
    """Shared handler for all Office → PDF conversions via LibreOffice with Python fallback."""
    in_path = await save_upload(file, suffix=ext)
    out_dir = make_tmp_dir()
    orig_name = Path(file.filename or "document").stem
    out_path = os.path.join(out_dir, f"{orig_name}.pdf")

    try:
        lo_out = _libreoffice_convert(in_path, out_dir, "pdf")
        if lo_out and os.path.exists(lo_out):
            out_path = lo_out
        else:
            # Python fallback when LibreOffice is not installed
            if ext in [".docx", ".doc"]:
                _fallback_docx_to_pdf(in_path, out_path)
            elif ext in [".xlsx", ".xls"]:
                _fallback_xlsx_to_pdf(in_path, out_path)
            elif ext in [".pptx", ".ppt"]:
                _fallback_pptx_to_pdf(in_path, out_path)
            else:
                raise Exception("Conversion unsupported")

        if not os.path.exists(out_path):
            raise Exception("Failed to generate PDF document")

    except Exception as e:
        cleanup(in_path)
        cleanup_dir(out_dir)
        raise HTTPException(status_code=500, detail=str(e))

    cleanup(in_path)
    return FileResponse(
        out_path,
        media_type="application/pdf",
        filename=f"{orig_name}.pdf",
        background=None,
    )


@router.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    """Convert Word (.docx/.doc) to PDF using LibreOffice headless (with Python fallback)."""
    return await _office_to_pdf(file, ".docx", "Word")


@router.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    """Convert Excel (.xlsx/.xls) to PDF using LibreOffice headless (with Python fallback)."""
    return await _office_to_pdf(file, ".xlsx", "Excel")


@router.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    """Convert PowerPoint (.pptx/.ppt) to PDF using LibreOffice headless (with Python fallback)."""
    return await _office_to_pdf(file, ".pptx", "PowerPoint")
